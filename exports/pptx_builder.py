"""PowerPoint export builder.

Implements CONTRACTS.md §5.3. Phase 5 supports modules ⊆
{demand, supply, financial, scorecard}; ``scenario`` raises
``NotImplementedError`` until Phase 7.

Default path is **programmatic** — slides are built from ``config.TKO`` design
tokens. ``template_path`` is an optional escape hatch: if provided (or if
``config.TKO_TEMPLATE_PATH`` exists on disk), that file is used as the slide
master and background painting is skipped so the template's design wins.

Charts are rendered to PNG in two layers:
1. **Plotly + kaleido** (SPEC §6.1) — preferred path; matches the on-screen
   look. Capped at ``_KALEIDO_TIMEOUT_S`` seconds per chart.
2. **Matplotlib (Agg)** via ``exports.chart_fallback`` — used when kaleido
   returns nothing within the timeout. Simpler visuals but renders entirely
   in-process (no chromium), so it's reliable on machines where kaleido's
   subprocess IPC is broken (see memory: kaleido-windows-hang).
3. Static "Chart render unavailable" placeholder — only reached if BOTH
   renderers fail. The underlying DataFrame is still attached to slide notes.
"""

from __future__ import annotations

import os
from concurrent.futures import ThreadPoolExecutor, TimeoutError as _FutureTimeout
from datetime import date, datetime
from io import BytesIO

import duckdb
import pandas as pd
import plotly.graph_objects as go
import plotly.io as pio
import streamlit as st
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

import config
from data import query
from data.query import Filters
from exports import chart_fallback

# Kaleido scope tweaks — disable mathjax fetch and pass chromium hardening
# flags. Kaleido 0.2.1 + Windows 11 has a known issue where the chromium
# network service crashes during chart-to-PNG conversion and IPC hangs;
# these flags reduce that risk but don't eliminate it. The hard cap is the
# per-chart timeout in ``_KALEIDO_TIMEOUT_S`` below.
try:
    pio.kaleido.scope.mathjax = None
    pio.kaleido.scope.chromium_args = (
        "--disable-gpu",
        "--no-sandbox",
        "--disable-dev-shm-usage",
        "--disable-software-rasterizer",
        "--single-process",
    )
except Exception:  # noqa: BLE001 — defensive; older kaleido versions
    pass

# Per-chart render timeout in seconds. First call cold-starts chromium so
# allow some headroom. If a chart times out, we fall back to matplotlib.
_KALEIDO_TIMEOUT_S = 20

# Build-scope circuit breaker. Once kaleido fails once during a build, skip it
# for subsequent charts in the same build and go straight to matplotlib. This
# keeps total build time tight on machines where kaleido is broken (~6s for the
# 6-chart bundle via matplotlib, vs 33s if we attempt kaleido on every chart).
# Reset at the top of every ``build()`` call so cache misses retry kaleido.
_kaleido_broken_this_build: bool = False

# ---------------------------------------------------------------------------
# Module support — declared up front so the dispatch matrix is visible
# ---------------------------------------------------------------------------

_SUPPORTED_MODULES = {"demand", "supply", "financial", "scorecard", "scenario"}

# Default scenario adjustments — used if the caller requested 'scenario' but
# didn't pass an adjustments dict (e.g. bundled export when sliders weren't
# touched). SPEC §5.5 — % volume adjustment vs consensus.
_DEFAULT_SCENARIO_ADJUSTMENTS = {"base": 0.0, "upside": 5.0, "downside": -5.0}

# Slide dimensions (10 × 5.625 in, 16:9)
_SLIDE_W = Inches(10)
_SLIDE_H = Inches(5.625)


# ---------------------------------------------------------------------------
# Color + style helpers
# ---------------------------------------------------------------------------


def _rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


_COLOR_VOID        = _rgb(config.TKO["colors"]["void"])
_COLOR_GRAPHITE    = _rgb(config.TKO["colors"]["graphite"])
_COLOR_COLD_SLATE  = _rgb(config.TKO["colors"]["cold_slate"])
_COLOR_PLASMA      = _rgb(config.TKO["colors"]["plasma"])
_COLOR_ACID        = _rgb(config.TKO["colors"]["acid"])
_COLOR_ICE_WHITE   = _rgb(config.TKO["colors"]["ice_white"])
_COLOR_ICE_DIM     = _rgb(config.TKO["colors"]["ice_dim"])
_COLOR_MIST        = _rgb(config.TKO["colors"]["mist"])

_FONT_DISPLAY = config.TKO["fonts"]["display"]
_FONT_BODY    = config.TKO["fonts"]["body"]

_RAG_RGB = {
    rag: _rgb(color) for rag, color in config.RAG_COLORS.items()
}


# ---------------------------------------------------------------------------
# Public entry
# ---------------------------------------------------------------------------


@st.cache_data(ttl=600, show_spinner=False)
def build(
    modules: list[str],
    filters: Filters,
    scenario_adjustments: dict | None = None,
    template_path: str | None = None,
) -> tuple[bytes, str]:
    """Render the requested modules to a ``.pptx`` byte stream.

    ``scenario`` raises ``NotImplementedError`` until Phase 7 — caller is
    expected to omit it from ``modules`` in Phase 5.
    """
    unsupported = [m for m in modules if m not in _SUPPORTED_MODULES]
    if unsupported:
        raise ValueError(
            f"pptx_builder: unknown modules {unsupported}. "
            f"Supported: {sorted(_SUPPORTED_MODULES)}."
        )
    if not modules:
        raise ValueError(
            f"modules must include at least one of {sorted(_SUPPORTED_MODULES)}"
        )

    cycle_month = _resolve_cycle_month(filters)
    use_template = _resolve_template_path(template_path)

    # Reset the kaleido circuit breaker per build — a previously cached run's
    # state would otherwise leak across calls.
    global _kaleido_broken_this_build
    _kaleido_broken_this_build = False

    if use_template:
        prs = Presentation(use_template)
        # Force 16:9 dims to match programmatic mode.
        prs.slide_width = _SLIDE_W
        prs.slide_height = _SLIDE_H
        painted_bg = False
    else:
        prs = Presentation()
        prs.slide_width = _SLIDE_W
        prs.slide_height = _SLIDE_H
        painted_bg = True

    _add_cover_slide(prs, filters, cycle_month, painted_bg)
    if "scorecard" in modules:
        _add_scorecard_slide(prs, filters, painted_bg)
    if "demand" in modules:
        _add_demand_slide(prs, filters, painted_bg)
    if "supply" in modules:
        _add_supply_slide(prs, filters, painted_bg)
    if "financial" in modules:
        _add_financial_slide(prs, filters, painted_bg)
    if "scenario" in modules:
        _add_scenario_slide(
            prs, filters,
            scenario_adjustments or _DEFAULT_SCENARIO_ADJUSTMENTS,
            painted_bg,
        )
    _add_appendix_slide(prs, modules, filters, cycle_month, painted_bg)

    buf = BytesIO()
    prs.save(buf)
    filename = f"S&OP_Report_{cycle_month}.pptx"
    return buf.getvalue(), filename


# ---------------------------------------------------------------------------
# Cycle-month + template resolution (mirrors xlsx_builder)
# ---------------------------------------------------------------------------


def _resolve_cycle_month(filters: Filters) -> str:
    dv = query.current_data_version("demand")
    if dv is not None:
        c = duckdb.connect(config.DUCKDB_PATH, read_only=True)
        try:
            row = c.execute(
                "SELECT MAX(period_date) FROM fact_demand WHERE data_version = ?",
                [dv],
            ).fetchone()
        finally:
            c.close()
        if row and row[0] is not None:
            return pd.Timestamp(row[0]).strftime("%Y-%m")
    if filters.period_to is not None:
        return pd.Timestamp(filters.period_to).strftime("%Y-%m")
    return pd.Timestamp(date.today()).strftime("%Y-%m")


def _resolve_template_path(template_path: str | None) -> str | None:
    if template_path and os.path.exists(template_path):
        return template_path
    if config.TKO_TEMPLATE_PATH and os.path.exists(config.TKO_TEMPLATE_PATH):
        return config.TKO_TEMPLATE_PATH
    return None


# ---------------------------------------------------------------------------
# Slide primitives — programmatic TKO build
# ---------------------------------------------------------------------------


def _blank_slide(prs: Presentation, painted_bg: bool):
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(layout)
    if painted_bg:
        # Paint a void rectangle across the entire slide as the background.
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, _SLIDE_W, _SLIDE_H,
        )
        bg.line.fill.background()
        bg.fill.solid()
        bg.fill.fore_color.rgb = _COLOR_VOID
        # Move to the back so subsequent shapes render on top.
        spTree = bg._element.getparent()
        spTree.remove(bg._element)
        spTree.insert(2, bg._element)
        # Plasma → acid accent strip at top.
        bar1 = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, Inches(7), Emu(45720),  # 0.05"
        )
        bar1.line.fill.background()
        bar1.fill.solid()
        bar1.fill.fore_color.rgb = _COLOR_PLASMA
        bar2 = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(7), 0, Inches(3), Emu(45720),
        )
        bar2.line.fill.background()
        bar2.fill.solid()
        bar2.fill.fore_color.rgb = _COLOR_ACID
    return slide


def _add_title(slide, text: str, *, top=Inches(0.2), left=Inches(0.4),
               width=Inches(9.2), height=Inches(0.6)) -> None:
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text.upper()
    run.font.name = _FONT_DISPLAY
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = _COLOR_ICE_WHITE


def _add_subtitle(slide, text: str, *, top=Inches(0.75), left=Inches(0.4),
                  width=Inches(9.2), height=Inches(0.3), color=None) -> None:
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = _FONT_BODY
    run.font.size = Pt(11)
    run.font.color.rgb = color or _COLOR_ICE_DIM


def _add_caption(slide, text: str, color=None) -> None:
    tb = slide.shapes.add_textbox(
        Inches(0.4), Inches(5.25), Inches(9.2), Inches(0.3),
    )
    p = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = _FONT_BODY
    run.font.size = Pt(9)
    run.font.color.rgb = color or _COLOR_MIST


def _fig_to_png(fig: go.Figure, *, width: int = 900, height: int = 540) -> bytes | None:
    """Render a Plotly figure to PNG via kaleido (SPEC §6.1: scale=2).

    Returns ``None`` if kaleido raises or doesn't return within
    ``_KALEIDO_TIMEOUT_S``. The caller falls back to matplotlib via
    ``chart_fallback.render``. The hung kaleido subprocess will eventually be
    cleaned up by the OS; we leak the ThreadPoolExecutor worker rather than
    block on it.

    Short-circuits to ``None`` if a previous call in this build already
    failed — see the circuit-breaker note on ``_kaleido_broken_this_build``.
    """
    global _kaleido_broken_this_build
    if _kaleido_broken_this_build:
        return None
    ex = ThreadPoolExecutor(max_workers=1)
    try:
        fut = ex.submit(
            fig.to_image, format="png", width=width, height=height, scale=2,
        )
        try:
            return fut.result(timeout=_KALEIDO_TIMEOUT_S)
        except (_FutureTimeout, Exception):  # noqa: BLE001
            _kaleido_broken_this_build = True
            return None
    finally:
        ex.shutdown(wait=False, cancel_futures=True)


def _add_image_from_fig(
    slide, fig: go.Figure, *, left, top, width, height,
    px_width: int = 900, px_height: int = 540,
    fallback_title: str = "Chart",
    fallback_kind: str | None = None,
    fallback_df: pd.DataFrame | None = None,
) -> None:
    """Render ``fig`` to PNG and place it on the slide.

    Three-tier fallback:
    1. Plotly + kaleido (preferred — visual parity with on-screen charts)
    2. Matplotlib via ``chart_fallback`` (when ``fallback_kind`` is provided)
    3. Static placeholder shape (only if both fail)
    """
    png = _fig_to_png(fig, width=px_width, height=px_height)
    if png is None and fallback_kind is not None:
        png = chart_fallback.render(
            fallback_kind, fallback_df,
            title=fallback_title, width_px=px_width, height_px=px_height,
        )
    if png is None:
        _add_chart_placeholder(
            slide, left=left, top=top, width=width, height=height,
            title=fallback_title,
        )
        return
    slide.shapes.add_picture(BytesIO(png), left, top, width=width, height=height)


def _add_chart_placeholder(slide, *, left, top, width, height, title: str) -> None:
    """Used when ``kaleido`` fails to render. Keeps the slide structurally
    complete; the underlying DataFrame is still attached to the slide notes.
    """
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _COLOR_GRAPHITE
    shape.line.color.rgb = _COLOR_COLD_SLATE
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.name = _FONT_DISPLAY
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = _COLOR_ICE_WHITE
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "Chart render unavailable — see slide notes / Excel for data."
    r2.font.name = _FONT_BODY
    r2.font.size = Pt(10)
    r2.font.color.rgb = _COLOR_MIST


def _attach_notes(slide, df: pd.DataFrame, *, label: str) -> None:
    """SPEC §6.1: dump the source DataFrame as CSV in the slide's notes pane."""
    notes = slide.notes_slide.notes_text_frame
    if not notes.text:
        notes.text = f"# {label}"
    else:
        p = notes.add_paragraph()
        p.text = f"# {label}"
    if df is None or df.empty:
        notes.add_paragraph().text = "(no data)"
        return
    csv = df.to_csv(index=False).strip()
    # python-pptx puts everything in one run; long strings are fine, but
    # truncate insanely large tables so the .pptx doesn't bloat.
    if len(csv) > 60_000:
        csv = csv[:60_000] + "\n…(truncated)"
    notes.add_paragraph().text = csv


# ---------------------------------------------------------------------------
# Cover
# ---------------------------------------------------------------------------


def _add_cover_slide(prs, f: Filters, cycle_month: str, painted_bg: bool) -> None:
    slide = _blank_slide(prs, painted_bg)

    _add_tko_wordmark(slide)

    # Big eyebrow + title block, centered vertically.
    eyebrow = slide.shapes.add_textbox(
        Inches(0.6), Inches(1.4), Inches(8.8), Inches(0.4),
    )
    p = eyebrow.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = f"S&OP CYCLE · {cycle_month}"
    run.font.name = _FONT_BODY
    run.font.size = Pt(12)
    run.font.color.rgb = _COLOR_ACID

    title = slide.shapes.add_textbox(
        Inches(0.6), Inches(1.9), Inches(8.8), Inches(1.2),
    )
    tp = title.text_frame.paragraphs[0]
    tr = tp.add_run()
    tr.text = "S&OP REVIEW"
    tr.font.name = _FONT_DISPLAY
    tr.font.size = Pt(54)
    tr.font.bold = True
    tr.font.color.rgb = _COLOR_ICE_WHITE

    meta = slide.shapes.add_textbox(
        Inches(0.6), Inches(3.2), Inches(8.8), Inches(0.4),
    )
    mp = meta.text_frame.paragraphs[0]
    mr = mp.add_run()
    mr.text = (
        f"Window {f.period_from:%Y-%m} → {f.period_to:%Y-%m} · "
        f"Generated {datetime.now():%Y-%m-%d %H:%M}"
    )
    mr.font.name = _FONT_BODY
    mr.font.size = Pt(13)
    mr.font.color.rgb = _COLOR_ICE_DIM

    # Filter summary block.
    rows = [
        ("Brands",     _filter_label(f.brands)),
        ("Categories", _filter_label(f.categories)),
        ("Channels",   _filter_label(f.channels)),
        ("Regions",    _filter_label(f.regions)),
        ("SKUs",       _filter_label(f.skus)),
    ]
    box = slide.shapes.add_textbox(
        Inches(0.6), Inches(3.85), Inches(8.8), Inches(1.3),
    )
    tf = box.text_frame
    for i, (label, value) in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        rl = p.add_run()
        rl.text = f"{label}: "
        rl.font.name = _FONT_BODY
        rl.font.size = Pt(10)
        rl.font.color.rgb = _COLOR_MIST
        rv = p.add_run()
        rv.text = value
        rv.font.name = _FONT_BODY
        rv.font.size = Pt(10)
        rv.font.color.rgb = _COLOR_ICE_DIM


def _filter_label(values: list[str]) -> str:
    return ", ".join(values) if values else "All"


def _add_tko_wordmark(slide) -> None:
    """Place the TKO wordmark in the bottom-right of the cover slide.

    PNG is rebuilt by ``assets/tko/build_wordmark.py``. If absent (asset zip not
    unzipped yet), the cover renders without it — no exception, no warning slide.
    """
    path = config.TKO_LOGO_PATH
    if not path or not os.path.exists(path):
        return
    # PNG aspect is ~3.25 : 1 (1040×320 native). Anchor flush to the bottom-right
    # margin so the wordmark sits beside the filter summary block on the cover.
    w = Inches(2.6)
    h = Inches(0.8)
    left = _SLIDE_W - w - Inches(0.4)
    top = _SLIDE_H - h - Inches(0.35)
    slide.shapes.add_picture(path, left, top, width=w, height=h)


# ---------------------------------------------------------------------------
# Scorecard slide
# ---------------------------------------------------------------------------


def _add_scorecard_slide(prs, f: Filters, painted_bg: bool) -> None:
    slide = _blank_slide(prs, painted_bg)
    _add_title(slide, "KPI Scorecard")
    _add_subtitle(slide, "Tier-1 KPIs with RAG status (SPEC §5.4)")

    df = query.scorecard_all(f)
    if df.empty:
        _add_caption(slide, "No data — upload demand, supply, or financial CSVs to populate the scorecard.")
        _attach_notes(slide, df, label="Scorecard")
        return

    headers = ["KPI", "Value", "RAG", "Thresholds"]
    rows = len(df) + 1
    cols = len(headers)
    table_shape = slide.shapes.add_table(
        rows, cols,
        Inches(0.4), Inches(1.2), Inches(9.2), Inches(3.6),
    )
    table = table_shape.table

    # Column widths — KPI / Value / RAG / Thresholds.
    table.columns[0].width = Inches(2.4)
    table.columns[1].width = Inches(1.6)
    table.columns[2].width = Inches(0.9)
    table.columns[3].width = Inches(4.3)

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _COLOR_COLD_SLATE
        _set_cell_text(cell, h, bold=True, size=11, color=_COLOR_ICE_WHITE)

    for i, r in df.iterrows():
        row_idx = i + 1
        _set_cell_text(
            table.cell(row_idx, 0), str(r["kpi_name"]),
            bold=True, size=11, color=_COLOR_ICE_WHITE, fill=_COLOR_GRAPHITE,
        )
        _set_cell_text(
            table.cell(row_idx, 1), _format_kpi_value(r),
            size=12, color=_COLOR_ICE_WHITE, fill=_COLOR_GRAPHITE,
        )
        rag = r["rag"] if pd.notna(r["rag"]) else "red"
        _set_cell_text(
            table.cell(row_idx, 2), rag.upper(),
            bold=True, size=11, color=_COLOR_VOID,
            fill=_RAG_RGB.get(rag, _RAG_RGB["red"]), align=PP_ALIGN.CENTER,
        )
        _set_cell_text(
            table.cell(row_idx, 3), str(r["threshold_explainer"]),
            size=10, color=_COLOR_ICE_DIM, fill=_COLOR_GRAPHITE,
        )

    counts = df["rag"].value_counts().to_dict()
    summary = (
        f"Green: {counts.get('green', 0)}  ·  "
        f"Amber: {counts.get('amber', 0)}  ·  "
        f"Red: {counts.get('red', 0)}"
    )
    _add_caption(slide, summary)
    _attach_notes(slide, df, label="Scorecard")


def _set_cell_text(cell, text: str, *, bold: bool = False, size: int = 11,
                   color: RGBColor | None = None, fill: RGBColor | None = None,
                   align: PP_ALIGN = PP_ALIGN.LEFT) -> None:
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    tf = cell.text_frame
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]
    p.alignment = align
    # Remove any default run text by clearing.
    p.text = ""
    run = p.add_run()
    run.text = text
    run.font.name = _FONT_BODY
    run.font.size = Pt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def _format_kpi_value(row: pd.Series) -> str:
    v = row["value"]
    if v is None or (isinstance(v, float) and pd.isna(v)):
        return "—"
    kpi, unit = row["kpi_name"], row["unit"]
    if kpi in ("Forecast Bias",):
        return f"{v:+.1f}%"
    if kpi in ("GM% vs Budget",):
        return f"{v:+.2f}pp"
    if unit == "%":
        return f"{v:.1f}%"
    if unit == "days":
        return f"{v:.1f} days"
    if unit == "pp":
        return f"{v:+.2f}pp"
    return f"{v:.2f}"


# ---------------------------------------------------------------------------
# Demand slide — Volume Bridge + Forecast vs Actuals
# ---------------------------------------------------------------------------


def _add_demand_slide(prs, f: Filters, painted_bg: bool) -> None:
    slide = _blank_slide(prs, painted_bg)
    _add_title(slide, "Demand Review")
    _add_subtitle(
        slide,
        f"Window {f.period_from:%Y-%m} → {f.period_to:%Y-%m} · Forecast accuracy, bias, mix",
    )

    bridge = query.demand_volume_bridge(f)
    series = query.demand_fcst_vs_actuals(f)

    fig_l = _fig_demand_bridge(bridge)
    fig_r = _fig_demand_fcst_vs_actuals(series)

    _add_image_from_fig(
        slide, fig_l,
        left=Inches(0.35), top=Inches(1.2),
        width=Inches(4.6), height=Inches(3.7),
        fallback_title="Volume Bridge",
        fallback_kind="demand_bridge", fallback_df=bridge,
    )
    _add_image_from_fig(
        slide, fig_r,
        left=Inches(5.05), top=Inches(1.2),
        width=Inches(4.6), height=Inches(3.7),
        fallback_title="Forecast vs Actuals",
        fallback_kind="demand_fcst_vs_actuals", fallback_df=series,
    )

    kpis = query.demand_kpis(f)
    _add_caption(
        slide,
        _safe(f"FA {kpis.fa_pct:.1f}%", kpis.fa_pct)
        + " · " + _safe(f"MAPE {kpis.mape:.1f}%", kpis.mape)
        + " · " + _safe(f"Bias {kpis.bias:+.1f}%", kpis.bias)
        + " · " + _safe(f"Volume {kpis.volume_total:,.0f} cases", kpis.volume_total, zero_ok=True),
    )
    _attach_notes(slide, bridge, label="Demand Volume Bridge")
    _attach_notes(slide, series, label="Demand Forecast vs Actuals")


def _fig_demand_bridge(df: pd.DataFrame) -> go.Figure:
    fig = go.Figure()
    if df.empty or df["value"].abs().sum() == 0:
        return _empty_fig("Volume Bridge")
    fig.add_trace(go.Waterfall(
        x=df["stage"], y=df["value"],
        measure=["absolute", "relative", "relative", "total"],
        text=[f"{v:,.0f}" for v in df["value"]],
        textposition="outside",
        connector={"line": {"color": config.TKO["colors"]["cold_slate"]}},
        increasing={"marker": {"color": config.SERIES["consensus_fcst"]}},
        decreasing={"marker": {"color": config.SERIES["statistical_fcst"]}},
        totals={"marker": {"color": config.SERIES["signal"]}},
    ))
    fig.update_layout(
        title="Volume Bridge",
        showlegend=False,
        yaxis_title="Cases",
        margin=dict(l=48, r=20, t=56, b=40),
    )
    return fig


def _fig_demand_fcst_vs_actuals(df: pd.DataFrame) -> go.Figure:
    fig = go.Figure()
    if df.empty:
        return _empty_fig("Forecast vs Actuals")
    fig.add_bar(
        x=df["period_date"], y=df["actuals"], name="Actuals",
        marker_color=config.SERIES["actuals"], opacity=0.85,
    )
    fig.add_scatter(
        x=df["period_date"], y=df["consensus_fcst"],
        name="Consensus", mode="lines+markers",
        line=dict(color=config.SERIES["consensus_fcst"], width=3),
    )
    fig.add_scatter(
        x=df["period_date"], y=df["statistical_fcst"],
        name="Statistical", mode="lines+markers",
        line=dict(color=config.SERIES["statistical_fcst"], width=2, dash="dash"),
    )
    fig.update_layout(
        title="Forecast vs Actuals",
        hovermode="x unified", yaxis_title="Cases",
        legend=dict(orientation="h", y=-0.18, x=0),
        margin=dict(l=48, r=20, t=56, b=60),
    )
    return fig


# ---------------------------------------------------------------------------
# Supply slide — DOS by SKU + Production Adherence
# ---------------------------------------------------------------------------


def _add_supply_slide(prs, f: Filters, painted_bg: bool) -> None:
    slide = _blank_slide(prs, painted_bg)
    _add_title(slide, "Supply Review")
    _add_subtitle(
        slide,
        f"Window {f.period_from:%Y-%m} → {f.period_to:%Y-%m} · Inventory cover and production",
    )

    dos_df = query.supply_dos_by_sku(f)
    adh_df = query.supply_production_adherence(f)
    fig_l = _fig_supply_dos(dos_df)
    fig_r = _fig_supply_adherence(adh_df)

    _add_image_from_fig(
        slide, fig_l,
        left=Inches(0.35), top=Inches(1.2),
        width=Inches(4.6), height=Inches(3.7),
        fallback_title="DOS by SKU",
        fallback_kind="supply_dos", fallback_df=dos_df,
    )
    _add_image_from_fig(
        slide, fig_r,
        left=Inches(5.05), top=Inches(1.2),
        width=Inches(4.6), height=Inches(3.7),
        fallback_title="Production Adherence",
        fallback_kind="supply_adherence", fallback_df=adh_df,
    )

    k = query.supply_kpis(f)
    _add_caption(
        slide,
        _safe(f"DOS {k.dos_avg:.1f} days", k.dos_avg)
        + " · " + _safe(f"Fill Rate {k.fill_rate:.1f}%", k.fill_rate)
        + " · " + _safe(f"Production Adherence {k.production_adherence:.1f}%", k.production_adherence),
    )
    _attach_notes(slide, dos_df, label="Supply DOS by SKU")
    _attach_notes(slide, adh_df, label="Supply Production Adherence")


def _fig_supply_dos(df: pd.DataFrame) -> go.Figure:
    if df.empty or df["dos"].isna().all():
        return _empty_fig("DOS by SKU")
    d = df.dropna(subset=["dos"]).sort_values("dos").tail(15)  # top 15 by DOS
    colors = [config.RAG_COLORS[r] for r in d["rag"]]
    fig = go.Figure(go.Bar(
        x=d["dos"], y=d["sku_name"], orientation="h",
        marker_color=colors,
        hovertemplate="%{y}<br>DOS: %{x:.1f} days<extra></extra>",
    ))
    fig.add_vline(
        x=config.RAG["dos"]["target_days"],
        line=dict(color=config.TKO["colors"]["acid"], width=2, dash="dot"),
    )
    fig.update_layout(
        title="DOS by SKU (latest period)",
        xaxis_title="Days of Supply",
        margin=dict(l=120, r=20, t=56, b=48),
    )
    return fig


def _fig_supply_adherence(df: pd.DataFrame) -> go.Figure:
    if df.empty:
        return _empty_fig("Production Adherence")
    fig = go.Figure()
    fig.add_bar(
        x=df["plant_name"], y=df["plan"], name="Plan",
        marker_color=config.SERIES["budget"],
    )
    fig.add_bar(
        x=df["plant_name"], y=df["actual"], name="Actual",
        marker_color=config.SERIES["consensus_fcst"],
    )
    fig.update_layout(
        title="Production Adherence by Plant",
        barmode="group", yaxis_title="Cases",
        legend=dict(orientation="h", y=-0.18, x=0),
        margin=dict(l=48, r=20, t=56, b=60),
    )
    return fig


# ---------------------------------------------------------------------------
# Financial slide — Revenue Waterfall + Revenue by Channel
# ---------------------------------------------------------------------------


def _add_financial_slide(prs, f: Filters, painted_bg: bool) -> None:
    slide = _blank_slide(prs, painted_bg)
    _add_title(slide, "Financial Review")
    _add_subtitle(
        slide,
        f"Window {f.period_from:%Y-%m} → {f.period_to:%Y-%m} · "
        f"Revenue vs Budget/LE · Currency: {config.CURRENCY}",
    )

    wf = query.financial_revenue_waterfall(f)
    chan = query.financial_by_channel(f)
    fig_l = _fig_financial_waterfall(wf)
    fig_r = _fig_financial_by_channel(chan)

    _add_image_from_fig(
        slide, fig_l,
        left=Inches(0.35), top=Inches(1.2),
        width=Inches(4.6), height=Inches(3.7),
        fallback_title="Revenue Waterfall",
        fallback_kind="financial_waterfall", fallback_df=wf,
    )
    _add_image_from_fig(
        slide, fig_r,
        left=Inches(5.05), top=Inches(1.2),
        width=Inches(4.6), height=Inches(3.7),
        fallback_title="Revenue by Channel",
        fallback_kind="financial_by_channel", fallback_df=chan,
    )

    k = query.financial_kpis(f)
    _add_caption(
        slide,
        _safe(f"Revenue {k.revenue_actual:,.0f} {config.CURRENCY}", k.revenue_actual, zero_ok=True)
        + " · " + _safe(f"vs Budget {k.revenue_vs_budget_pct:+.1f}%", k.revenue_vs_budget_pct)
        + " · " + _safe(f"vs LE {k.revenue_vs_le_pct:+.1f}%", k.revenue_vs_le_pct)
        + " · " + _safe(f"GM {k.gm_pct:.1f}%", k.gm_pct),
    )
    _attach_notes(slide, wf, label="Financial Revenue Waterfall")
    _attach_notes(slide, chan, label="Financial Revenue by Channel")


def _fig_financial_waterfall(df: pd.DataFrame) -> go.Figure:
    if df.empty or df["value"].abs().sum() == 0:
        return _empty_fig("Revenue Waterfall")
    budget = float(df.loc[df["stage"] == "Budget", "value"].iloc[0])
    le     = float(df.loc[df["stage"] == "LE",     "value"].iloc[0])
    actual = float(df.loc[df["stage"] == "Actuals", "value"].iloc[0])
    fig = go.Figure(go.Waterfall(
        x=["Budget", "Δ to LE", "Δ to Actual", "Actual"],
        y=[budget, le - budget, actual - le, actual],
        measure=["absolute", "relative", "relative", "total"],
        text=[
            f"{budget:,.0f}",
            f"{(le - budget):+,.0f}",
            f"{(actual - le):+,.0f}",
            f"{actual:,.0f}",
        ],
        textposition="outside",
        connector={"line": {"color": config.TKO["colors"]["cold_slate"]}},
        increasing={"marker": {"color": config.SERIES["consensus_fcst"]}},
        decreasing={"marker": {"color": config.SERIES["statistical_fcst"]}},
        totals={"marker": {"color": config.SERIES["signal"]}},
    ))
    fig.update_layout(
        title="Revenue Waterfall",
        showlegend=False,
        yaxis_title=f"Revenue ({config.CURRENCY})",
        margin=dict(l=64, r=20, t=56, b=40),
    )
    return fig


def _fig_financial_by_channel(df: pd.DataFrame) -> go.Figure:
    if df.empty or df["revenue_actual"].abs().sum() == 0:
        return _empty_fig("Revenue by Channel")
    fig = go.Figure(go.Treemap(
        labels=df["channel_name"],
        parents=[""] * len(df),
        values=df["revenue_actual"],
        marker=dict(
            colors=df["revenue_actual"],
            colorscale=[
                [0.0, config.TKO["colors"]["graphite"]],
                [1.0, config.TKO["colors"]["plasma"]],
            ],
        ),
        texttemplate="<b>%{label}</b><br>%{value:,.0f}<br>%{percentRoot:.1%}",
    ))
    fig.update_layout(
        title="Revenue by Channel",
        margin=dict(l=8, r=8, t=56, b=8),
    )
    return fig


# ---------------------------------------------------------------------------
# Scenario slide — 3-scenario comparison table + revenue bar chart
# ---------------------------------------------------------------------------


def _scenario_summary_df(f: Filters, adjustments: dict) -> pd.DataFrame:
    """Aggregate per-scenario (volume_cases, revenue) across (sku, channel).

    Mirrors xlsx_builder._scenario_summary but kept local to avoid the
    exports/ ↔ exports/ import (chart_fallback is the only cross-module import
    in this builder).
    """
    base = query.scenario_baseline(f)
    price = query.scenario_price_proxy(f)
    if base.empty:
        return pd.DataFrame(
            columns=["scenario", "adjustment_pct", "volume_cases", "revenue"]
        )
    merged = base.merge(price, on=["sku_code", "channel_code"], how="left")
    rows = []
    for name in ("Base", "Upside", "Downside"):
        adj = float(adjustments.get(name.lower(), 0.0))
        vol = merged["consensus_volume_cases"] * (1.0 + adj / 100.0)
        rev = (vol * merged["avg_price"]).fillna(0.0)
        rows.append({
            "scenario": name,
            "adjustment_pct": adj,
            "volume_cases": float(vol.sum()),
            "revenue": float(rev.sum()),
        })
    return pd.DataFrame(rows)


def _add_scenario_slide(prs, f: Filters, adjustments: dict, painted_bg: bool) -> None:
    slide = _blank_slide(prs, painted_bg)
    _add_title(slide, "Scenario Comparison")
    _add_subtitle(
        slide,
        f"Window {f.period_from:%Y-%m} → {f.period_to:%Y-%m} · "
        f"Volume × trailing-3M avg price · Currency: {config.CURRENCY}",
    )

    summary = _scenario_summary_df(f, adjustments)
    if summary.empty:
        _add_caption(
            slide,
            "No forward planning horizon — actuals extend through the window, "
            "so scenarios have no months to project.",
        )
        _attach_notes(slide, summary, label="Scenario Comparison (empty)")
        return

    # Table on the left.
    headers = ["Scenario", "% vs Consensus", "Volume (cases)",
               f"Revenue ({config.CURRENCY})"]
    table_shape = slide.shapes.add_table(
        len(summary) + 1, len(headers),
        Inches(0.35), Inches(1.2), Inches(5.2), Inches(2.4),
    )
    table = table_shape.table
    table.columns[0].width = Inches(1.2)
    table.columns[1].width = Inches(1.3)
    table.columns[2].width = Inches(1.3)
    table.columns[3].width = Inches(1.4)

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _COLOR_COLD_SLATE
        _set_cell_text(cell, h, bold=True, size=11, color=_COLOR_ICE_WHITE)

    palette = {
        "Base":     _COLOR_GRAPHITE,
        "Upside":   _RAG_RGB["green"],
        "Downside": _RAG_RGB["red"],
    }
    for i, r in summary.iterrows():
        bg = palette.get(r["scenario"], _COLOR_GRAPHITE)
        text_color = (
            _COLOR_VOID if r["scenario"] in ("Upside", "Downside")
            else _COLOR_ICE_WHITE
        )
        _set_cell_text(
            table.cell(i + 1, 0), str(r["scenario"]),
            bold=True, size=11, color=text_color, fill=bg,
        )
        _set_cell_text(
            table.cell(i + 1, 1), f"{r['adjustment_pct']:+.1f}%",
            size=11, color=text_color, fill=bg, align=PP_ALIGN.RIGHT,
        )
        _set_cell_text(
            table.cell(i + 1, 2), f"{r['volume_cases']:,.0f}",
            size=11, color=text_color, fill=bg, align=PP_ALIGN.RIGHT,
        )
        _set_cell_text(
            table.cell(i + 1, 3), f"{r['revenue']:,.0f}",
            bold=True, size=11, color=text_color, fill=bg, align=PP_ALIGN.RIGHT,
        )

    # Chart on the right.
    fig = _fig_scenario_revenue(summary)
    _add_image_from_fig(
        slide, fig,
        left=Inches(5.65), top=Inches(1.2),
        width=Inches(4.0), height=Inches(3.6),
        fallback_title="Scenario Revenue",
        fallback_kind="scenario_revenue", fallback_df=summary,
    )

    base_rev = float(summary.loc[summary["scenario"] == "Base", "revenue"].iloc[0])
    up_rev   = float(summary.loc[summary["scenario"] == "Upside", "revenue"].iloc[0])
    down_rev = float(summary.loc[summary["scenario"] == "Downside", "revenue"].iloc[0])

    def _delta(v: float) -> str:
        if base_rev == 0:
            return "—"
        return f"{(v - base_rev) / base_rev * 100:+.1f}%"

    _add_caption(
        slide,
        f"Base {base_rev:,.0f} · "
        f"Upside {up_rev:,.0f} ({_delta(up_rev)}) · "
        f"Downside {down_rev:,.0f} ({_delta(down_rev)})",
    )
    _attach_notes(slide, summary, label="Scenario Comparison")


def _fig_scenario_revenue(df: pd.DataFrame) -> go.Figure:
    if df.empty:
        return _empty_fig("Scenario Revenue")
    palette = {
        "Base":     config.SERIES["consensus_fcst"],
        "Upside":   config.RAG_COLORS["green"],
        "Downside": config.RAG_COLORS["red"],
    }
    colors = [palette.get(s, config.SERIES["consensus_fcst"]) for s in df["scenario"]]
    fig = go.Figure(go.Bar(
        x=df["scenario"], y=df["revenue"],
        marker_color=colors,
        text=[f"{v:,.0f}" for v in df["revenue"]],
        textposition="outside",
    ))
    fig.update_layout(
        title="Scenario Revenue",
        showlegend=False,
        yaxis_title=f"Revenue ({config.CURRENCY})",
        margin=dict(l=56, r=24, t=56, b=40),
    )
    return fig


# ---------------------------------------------------------------------------
# Appendix
# ---------------------------------------------------------------------------


def _add_appendix_slide(
    prs, modules: list[str], f: Filters, cycle_month: str, painted_bg: bool,
) -> None:
    slide = _blank_slide(prs, painted_bg)
    _add_title(slide, "Appendix · Data & Scope")

    lines = [
        ("Cycle month",     cycle_month),
        ("Window",          f"{f.period_from:%Y-%m} → {f.period_to:%Y-%m}"),
        ("Modules",         ", ".join(modules)),
        ("Currency",        config.CURRENCY),
        ("Demand version",    _version_label("demand")),
        ("Supply version",    _version_label("supply")),
        ("Financial version", _version_label("financial")),
        ("Brands",     _filter_label(f.brands)),
        ("Categories", _filter_label(f.categories)),
        ("Channels",   _filter_label(f.channels)),
        ("Regions",    _filter_label(f.regions)),
        ("SKUs",       _filter_label(f.skus)),
    ]
    box = slide.shapes.add_textbox(
        Inches(0.6), Inches(1.2), Inches(8.8), Inches(3.8),
    )
    tf = box.text_frame
    for i, (label, value) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        rl = p.add_run()
        rl.text = f"{label}: "
        rl.font.name = _FONT_BODY
        rl.font.size = Pt(11)
        rl.font.color.rgb = _COLOR_MIST
        rv = p.add_run()
        rv.text = value
        rv.font.name = _FONT_BODY
        rv.font.size = Pt(11)
        rv.font.color.rgb = _COLOR_ICE_DIM

    _add_caption(slide, "Source CSVs are retained in storage/raw/ for audit.")


def _version_label(domain: str) -> str:
    v = query.current_data_version(domain)
    return f"v{v}" if v else "—"


# ---------------------------------------------------------------------------
# Shared utilities
# ---------------------------------------------------------------------------


def _empty_fig(title: str) -> go.Figure:
    fig = go.Figure()
    fig.add_annotation(
        text="No data in window",
        showarrow=False,
        font=dict(color=config.TKO["colors"]["mist"], size=14),
    )
    fig.update_layout(
        title=title,
        xaxis=dict(visible=False), yaxis=dict(visible=False),
        margin=dict(l=20, r=20, t=56, b=20),
    )
    return fig


def _safe(text: str, v: float | None, *, zero_ok: bool = False) -> str:
    """Render a metric chip text, or '—' for NaN. ``zero_ok`` lets 0.0 render."""
    if v is None:
        return text.split(" ")[0] + " —"
    if isinstance(v, float) and pd.isna(v):
        return text.split(" ")[0] + " —"
    if not zero_ok and v == 0:
        # 0 from "no rows" looks identical to a real zero — show "—".
        return text.split(" ")[0] + " —"
    return text
